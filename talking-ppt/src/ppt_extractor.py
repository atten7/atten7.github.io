import os
from pptx import Presentation
from PIL import Image
from io import BytesIO
import numpy as np
from transformers import BlipProcessor, BlipForConditionalGeneration

# BLIP 모델 로드
processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")

# 단색 이미지 제외
def is_solid_color(image_bytes, tolerance=10):
    try:
        pil_image = Image.open(BytesIO(image_bytes)).convert("RGB")
        img_array = np.array(pil_image)
        
        if img_array.shape[0] < 5 or img_array.shape[1] < 5:
            return False

        mean_color = np.mean(img_array, axis=(0, 1))
        std_dev = np.std(img_array.reshape(-1, 3), axis=0)

        return np.all(std_dev < tolerance)

    except Exception as e:
        print(f"[오류] 이미지 처리: {e}")
        return False

# 텍스트, 메모, 이미지(캡션) 추출
def extract_ppt_content(ppt_name):
    try:
        prs = Presentation(ppt_name)
    except FileNotFoundError:
        print(f"[오류] 파일 '{ppt_name}' 존재 안함")
        return []

    slides_data = []
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide_area = slide_width * slide_height

    for idx, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "slide_number": idx,
            "text": [],
            "notes": "",
            "images": []
        }

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_info["text"].append(shape.text.strip())
        
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text if notes_slide and notes_slide.notes_text_frame else ""
        slide_info["notes"] = notes_text.strip()
        
        for shape in slide.shapes:
            if shape.shape_type == 13: #사진
                image = shape.image
                image_bytes = image.blob
                
                if is_solid_color(image_bytes):
                    print(f"[로그] 슬라이드 {idx} - 단색 이미지 제외")
                    continue
                
                image_area = shape.width * shape.height
                if image_area / slide_area < 0.01:
                    print(f"[로그] 슬라이드 {idx} - 작은 이미지 제외")
                    continue

                pil_image = Image.open(BytesIO(image_bytes)).convert("RGB")
                inputs = processor(pil_image, return_tensors="pt")
                out = blip_model.generate(**inputs)
                caption = processor.decode(out[0], skip_special_tokens=True)
                
                if any(keyword in caption for keyword in ["logo", "background", "template", "design", "diagram", "chart"]):
                    print(f"[로그] 슬라이드 {idx} - 키워드 '{caption}' 이미지 제외")
                    continue
                
                slide_info["images"].append(caption)
        
        slides_data.append(slide_info)

    return slides_data